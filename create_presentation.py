#!/usr/bin/env python3
# -*- coding: utf-8 -*-

from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.dml.color import RGBColor

def create_presentation():
    """TS GenAI Gateway 제품소개서 생성"""

    prs = Presentation()
    prs.slide_width = Inches(10)
    prs.slide_height = Inches(7.5)

    # 색상 팔레트 (Deep Blue + Electric Cyan)
    PRIMARY_COLOR = RGBColor(20, 52, 164)      # Deep Blue
    ACCENT_COLOR = RGBColor(0, 191, 255)       # Electric Cyan
    TEXT_COLOR = RGBColor(33, 33, 33)          # Dark Gray
    LIGHT_BG = RGBColor(245, 248, 250)         # Light Blue Gray

    def add_title_slide(title, subtitle):
        """타이틀 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout

        # 배경색
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = PRIMARY_COLOR

        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(2.5), Inches(9), Inches(1.5)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.alignment = PP_ALIGN.CENTER
        title_para.font.size = Pt(60)
        title_para.font.bold = True
        title_para.font.color.rgb = RGBColor(255, 255, 255)

        # 서브타이틀
        subtitle_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(4.2), Inches(9), Inches(1)
        )
        subtitle_frame = subtitle_box.text_frame
        subtitle_frame.text = subtitle
        subtitle_para = subtitle_frame.paragraphs[0]
        subtitle_para.alignment = PP_ALIGN.CENTER
        subtitle_para.font.size = Pt(28)
        subtitle_para.font.color.rgb = ACCENT_COLOR

    def add_content_slide(title, content_items, style="bullet"):
        """콘텐츠 슬라이드 추가"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 콘텐츠
        content_box = slide.shapes.add_textbox(
            Inches(0.8), Inches(1.8), Inches(8.4), Inches(5)
        )
        text_frame = content_box.text_frame
        text_frame.word_wrap = True

        for i, item in enumerate(content_items):
            if i > 0:
                text_frame.add_paragraph()
            p = text_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(20)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(14)
            p.level = 0

    def add_two_column_slide(title, left_items, right_items):
        """2단 레이아웃 슬라이드"""
        slide = prs.slides.add_slide(prs.slide_layouts[6])

        # 타이틀
        title_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(0.5), Inches(9), Inches(0.8)
        )
        title_frame = title_box.text_frame
        title_frame.text = title
        title_para = title_frame.paragraphs[0]
        title_para.font.size = Pt(44)
        title_para.font.bold = True
        title_para.font.color.rgb = PRIMARY_COLOR

        # 왼쪽 콘텐츠
        left_box = slide.shapes.add_textbox(
            Inches(0.5), Inches(1.8), Inches(4.2), Inches(5)
        )
        left_frame = left_box.text_frame
        for i, item in enumerate(left_items):
            if i > 0:
                left_frame.add_paragraph()
            p = left_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(12)

        # 오른쪽 콘텐츠
        right_box = slide.shapes.add_textbox(
            Inches(5.3), Inches(1.8), Inches(4.2), Inches(5)
        )
        right_frame = right_box.text_frame
        for i, item in enumerate(right_items):
            if i > 0:
                right_frame.add_paragraph()
            p = right_frame.paragraphs[i]
            p.text = item
            p.font.size = Pt(18)
            p.font.color.rgb = TEXT_COLOR
            p.space_after = Pt(12)

    # 슬라이드 1: 타이틀
    add_title_slide(
        "TS GenAI Gateway",
        "모든 AI를 하나로, 사용한 만큼만"
    )

    # 슬라이드 2: 문제 제기
    add_content_slide(
        "기업이 AI 도입하며 겪는 현실",
        [
            "💸 AI 구독료 폭발",
            "   • ChatGPT Plus $20 + Claude Pro $20 + Gemini $20 = 팀 10명이면 월 $600",
            "   • 실제 사용량은 50%인데 100% 요금 지불",
            "",
            "🔒 보안 리스크",
            "   • 직원들이 개인 계정으로 회사 기밀 입력",
            "   • 데이터가 어디로 가는지 추적 불가",
            "",
            "🎯 관리 불가",
            "   • 어떤 팀이 얼마나 사용하는지 모름",
            "   • 프롬프트 품질 천차만별",
            "   • 노하우가 개인에게만 축적"
        ]
    )

    # 슬라이드 3: 기존 해결책의 한계
    add_content_slide(
        "기존 해결책이 실패하는 이유",
        [
            "❌ 개인 계정 사용",
            "   → 보안 위험, 비용 통제 불가, 지식 공유 안 됨",
            "",
            "❌ 한 가지 AI만 선택",
            "   → 용도별 최적 모델 활용 불가 (GPT-4 vs Claude vs Gemini)",
            "",
            "❌ 자체 개발",
            "   → 개발 비용 $50K+, 유지보수 부담, 6개월 지연",
            "",
            "✅ 필요한 것",
            "   → 통합 접근 + 실시간 비용 관리 + 엔터프라이즈 보안"
        ]
    )

    # 슬라이드 4: 솔루션 - TS GenAI Gateway
    add_content_slide(
        "TS GenAI Gateway: 모든 AI를 하나로",
        [
            "🎯 핵심 가치",
            "",
            "1️⃣ 하나의 인터페이스로 10개 AI 모델 접근",
            "   • ChatGPT, Claude, Gemini, Llama 등 원하는 모델 즉시 전환",
            "",
            "2️⃣ 토큰 단위 실시간 비용 추적",
            "   • 이 대화가 지금 얼마? → 즉시 확인",
            "   • 예산 80% 도달 시 자동 알림",
            "",
            "3️⃣ 엔터프라이즈 보안 및 거버넌스",
            "   • 민감정보 자동 필터링 (Presidio SDK)",
            "   • 모든 대화 감사 로그",
            "",
            "4️⃣ 팀 협업 및 지식 축적",
            "   • 프롬프트 템플릿 200+ 제공 및 팀 공유",
            "   • 대화 히스토리 검색 및 재사용"
        ]
    )

    # 슬라이드 5: 핵심 기능 1 - 통합 AI 접근
    add_two_column_slide(
        "핵심 기능 1: 통합 AI 모델 허브",
        [
            "🤖 지원 모델",
            "• OpenAI GPT-4, GPT-3.5",
            "• Anthropic Claude 3 Opus/Sonnet",
            "• Google Gemini Pro/Ultra",
            "• Meta Llama 3",
            "• 오픈소스 모델 (Mistral 등)",
            "",
            "🔄 실시간 전환",
            "• 대화 중 모델 변경 가능",
            "• 용도별 최적 모델 자동 추천",
            "• A/B 테스트로 성능 비교"
        ],
        [
            "📊 모델 카탈로그",
            "• 각 모델의 상세 스펙",
            "  - 컨텍스트 윈도우",
            "  - 토큰당 가격",
            "  - 지원 기능 (이미지, 파일 등)",
            "",
            "🎯 사용 사례",
            "• 코드 생성 → GPT-4",
            "• 긴 문서 분석 → Claude",
            "• 빠른 답변 → Gemini",
            "• 비용 절감 → Llama 3"
        ]
    )

    # 슬라이드 6: 핵심 기능 2 - 비용 관리
    add_content_slide(
        "핵심 기능 2: 스마트한 비용 관리",
        [
            "💰 실시간 토큰 추적",
            "   • Input/Output 토큰 정확히 계산",
            "   • 대화 중 실시간 비용 표시",
            "   • 일/주/월별 사용량 분석 대시보드",
            "",
            "🎯 유연한 한도 설정",
            "   • 사용자별, 팀별, 프로젝트별 예산 설정",
            "   • 소프트 리밋 (알림) vs 하드 리밋 (차단)",
            "   • 80% 도달 시 이메일/슬랙 알림",
            "",
            "📊 예측 및 최적화",
            "   • 향후 비용 추이 예측",
            "   • 비효율적 사용 패턴 자동 감지",
            "   • 월말 청구서 대신 실시간 투명성"
        ]
    )

    # 슬라이드 7: 핵심 기능 3 - 보안
    add_content_slide(
        "핵심 기능 3: 엔터프라이즈급 보안",
        [
            "🔐 민감정보 보호 (Presidio SDK)",
            "   • 이메일, 전화번호, 주민번호 자동 마스킹",
            "   • 전송 전 필터링 → AI에 도달하기 전 차단",
            "",
            "📝 완전한 감사 추적",
            "   • 모든 대화, 파일 업로드 기록",
            "   • 누가, 언제, 무엇을, 어느 모델에",
            "   • 규정 준수 (GDPR, HIPAA 대응)",
            "",
            "🔑 세밀한 접근 제어",
            "   • OAuth 기반 소셜 로그인",
            "   • 역할 기반 권한 관리 (Admin/User)",
            "   • API 키 발급 및 관리",
            "",
            "🛡️ 보안 이벤트 모니터링",
            "   • 비정상 로그인 시도 자동 탐지",
            "   • 의심스러운 프롬프트 패턴 알림"
        ]
    )

    # 슬라이드 8: 핵심 기능 4 - 협업
    add_two_column_slide(
        "핵심 기능 4: 협업과 생산성",
        [
            "👥 팀 협업",
            "• 그룹 채팅 (실시간 동기화)",
            "• 대화 공유 및 컨텍스트 유지",
            "• 실시간 다국어 번역",
            "",
            "📚 지식 관리",
            "• 프롬프트 템플릿 200+ 제공",
            "  - 법률, 마케팅, 개발 등",
            "• 커스텀 템플릿 저장/공유",
            "• Git 스타일 버전 관리",
            "",
            "🔍 스마트 검색",
            "• 전체 대화 히스토리 검색",
            "• AI 기반 자동 요약",
            "• 북마크 및 태그 분류"
        ],
        [
            "🗂️ 프로젝트 관리",
            "• 프로젝트별 컨텍스트 분리",
            "• 접근 권한 설정",
            "  (비공개/팀/조직)",
            "• 프로젝트별 비용 추적",
            "",
            "📎 파일 지원",
            "• 문서/이미지 업로드",
            "• 버전 관리 (v1.0, v1.1)",
            "• PDF, PNG, CSV 등 지원",
            "",
            "🎨 사용자 경험",
            "• 스트리밍 응답 (실시간)",
            "• 응답 중단 기능",
            "• 모바일 최적화"
        ]
    )

    # 슬라이드 9: 차별점
    add_content_slide(
        "왜 TS GenAI Gateway인가?",
        [
            "⚡ 유일한 토큰 단위 실시간 추적",
            "   • 경쟁사: 월말 청구서 보고 놀람",
            "   • 우리: 매 순간 정확히 파악, 예산 초과 Zero",
            "",
            "🔄 진정한 멀티 모델",
            "   • 경쟁사: 한 모델만 또는 복잡한 전환",
            "   • 우리: 클릭 한 번에 10개 모델 자유롭게",
            "",
            "🛡️ 내장된 엔터프라이즈 보안",
            "   • 경쟁사: 보안은 별도 구축",
            "   • 우리: Presidio 기반 민감정보 자동 필터링 내장",
            "",
            "📚 검증된 템플릿 200+",
            "   • 경쟁사: 빈 화면에서 시작",
            "   • 우리: 업종별 베스트 프랙티스 즉시 활용",
            "",
            "💰 사용한 만큼만 (Pay-as-you-go)",
            "   • 경쟁사: 월 $20 고정 요금",
            "   • 우리: 실제 사용량만 결제, 평균 70% 절감"
        ]
    )

    # 슬라이드 10: 주요 기능 요약
    add_two_column_slide(
        "Phase 1 주요 기능 (28개 기능)",
        [
            "👤 사용자 관리",
            "• 개인 프로필 및 선호 설정",
            "• 프롬프트 템플릿 관리",
            "• 토큰 사용량 추적/한도",
            "• 대화 히스토리 관리",
            "• 그룹 채팅 및 실시간 번역",
            "",
            "🔐 인증 및 보안",
            "• OAuth 소셜 로그인",
            "• API 키 관리",
            "• 보안 이벤트 로깅",
            "• 민감정보 필터 (Presidio)",
            "",
            "📁 프로젝트 관리",
            "• 프로젝트 생성 및 설정",
            "• 태그/메타데이터 관리"
        ],
        [
            "🤖 AI 모델 통합",
            "• 10+ 모델 카탈로그",
            "• 호환 API 및 어댑터",
            "• OpenAI 호환 엔드포인트",
            "",
            "💬 대화 관리",
            "• 컨텍스트 유지",
            "• 스트리밍 응답 (SSE)",
            "• 북마크 및 태그",
            "",
            "💰 크레딧 시스템",
            "• 종량제 결제",
            "• 사용량 한도 및 알림",
            "",
            "📊 통계/모니터링",
            "• 사용량 대시보드",
            "• 비용 분석 및 예측"
        ]
    )

    # 슬라이드 11: 시작하기 (CTA)
    add_content_slide(
        "지금 바로 시작하세요",
        [
            "🚀 무료 체험",
            "   • 신규 가입 시 $10 크레딧 제공 (30일 유효)",
            "   • 신용카드 등록 불필요",
            "   • 5분이면 설정 완료",
            "",
            "💳 유연한 요금제",
            "   • Pay-as-you-go: 사용한 만큼만",
            "   • 팀 플랜: 중앙 관리 + 할인",
            "   • 엔터프라이즈: 전용 지원 + SLA",
            "",
            "🤝 지원",
            "   • 24/7 이메일 지원",
            "   • 라이브 채팅",
            "   • 온보딩 가이드 및 튜토리얼",
            "",
            "📧 문의",
            "   • 이메일: contact@tsgenai.com",
            "   • 웹사이트: www.tsgenai.com",
            "   • 데모 요청: 즉시 예약 가능"
        ]
    )

    # 슬라이드 12: 감사 슬라이드
    slide = prs.slides.add_slide(prs.slide_layouts[6])
    background = slide.background
    fill = background.fill
    fill.solid()
    fill.fore_color.rgb = PRIMARY_COLOR

    thanks_box = slide.shapes.add_textbox(
        Inches(0.5), Inches(2.5), Inches(9), Inches(2)
    )
    thanks_frame = thanks_box.text_frame
    thanks_frame.text = "감사합니다\n\nTS GenAI Gateway\n모든 AI를 하나로, 사용한 만큼만"

    for paragraph in thanks_frame.paragraphs:
        paragraph.alignment = PP_ALIGN.CENTER
        paragraph.font.size = Pt(40)
        paragraph.font.color.rgb = RGBColor(255, 255, 255)

    thanks_frame.paragraphs[0].font.size = Pt(48)
    thanks_frame.paragraphs[0].font.bold = True

    # 저장
    filename = "TS_GenAI_Gateway_제품소개서.pptx"
    prs.save(filename)
    print(f"✅ 제품소개서가 생성되었습니다: {filename}")
    return filename

if __name__ == "__main__":
    create_presentation()
